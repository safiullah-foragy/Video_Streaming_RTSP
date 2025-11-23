"""
Create a professional PowerPoint presentation for the Video Streaming Server project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    dark_bg = RGBColor(18, 18, 18)
    accent_red = RGBColor(239, 68, 68)
    white = RGBColor(255, 255, 255)
    gray = RGBColor(200, 200, 200)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = dark_bg
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title = title_frame.add_paragraph()
    title.text = "🎬 Video Streaming Server"
    title.font.size = Pt(54)
    title.font.bold = True
    title.font.color.rgb = white
    title.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle = subtitle_frame.add_paragraph()
    subtitle.text = "RTSP/RTP Protocol & YouTube-Style Web Interface"
    subtitle.font.size = Pt(24)
    subtitle.font.color.rgb = accent_red
    subtitle.alignment = PP_ALIGN.CENTER
    
    author_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    author_frame = author_box.text_frame
    author = author_frame.add_paragraph()
    author.text = "Safiullah Foragy"
    author.font.size = Pt(18)
    author.font.color.rgb = gray
    author.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = dark_bg
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title = title_frame.add_paragraph()
    title.text = "📋 Project Overview"
    title.font.size = Pt(44)
    title.font.bold = True
    title.font.color.rgb = accent_red
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    points = [
        "Comprehensive video streaming solution",
        "Dual interface: Web-based + RTSP/RTP protocol",
        "YouTube-inspired modern UI with dark theme",
        "Custom video player with advanced controls",
        "Network streaming across devices",
        "Production-ready and deployed live"
    ]
    
    for point in points:
        p = content_frame.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(24)
        p.font.color.rgb = white
        p.space_before = Pt(12)
        p.level = 0
    
    # Slide 3: Key Features - Web Interface
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = dark_bg
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title = title_frame.add_paragraph()
    title.text = "🌐 Web Interface Features"
    title.font.size = Pt(44)
    title.font.bold = True
    title.font.color.rgb = accent_red
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    features = [
        "🎨 YouTube-Style UI - Modern, responsive dark theme",
        "🎥 Custom Video Player - Professional controls",
        "⏩ 2-Second Seek Buttons - Quick navigation",
        "🔊 Smart Audio Handling - Auto-unmute on interaction",
        "⌨️ Keyboard Shortcuts - Full keyboard control",
        "📱 Responsive Design - Works on all devices"
    ]
    
    for feature in features:
        p = content_frame.add_paragraph()
        p.text = feature
        p.font.size = Pt(22)
        p.font.color.rgb = white
        p.space_before = Pt(14)
    
    # Slide 4: Technical Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = dark_bg
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title = title_frame.add_paragraph()
    title.text = "🛠️ Technology Stack"
    title.font.size = Pt(44)
    title.font.bold = True
    title.font.color.rgb = accent_red
    
    # Backend
    backend_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4), Inches(2.5))
    backend_frame = backend_box.text_frame
    backend_frame.word_wrap = True
    
    b_title = backend_frame.add_paragraph()
    b_title.text = "Backend"
    b_title.font.size = Pt(28)
    b_title.font.bold = True
    b_title.font.color.rgb = accent_red
    
    backend_items = ["Python 3.8+", "Flask 3.1.2", "OpenCV", "Gunicorn"]
    for item in backend_items:
        p = backend_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(20)
        p.font.color.rgb = white
        p.space_before = Pt(8)
    
    # Frontend
    frontend_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4), Inches(2.5))
    frontend_frame = frontend_box.text_frame
    frontend_frame.word_wrap = True
    
    f_title = frontend_frame.add_paragraph()
    f_title.text = "Frontend"
    f_title.font.size = Pt(28)
    f_title.font.bold = True
    f_title.font.color.rgb = accent_red
    
    frontend_items = ["HTML5 Video API", "CSS3 Dark Theme", "Vanilla JavaScript", "Font Awesome Icons"]
    for item in frontend_items:
        p = frontend_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(20)
        p.font.color.rgb = white
        p.space_before = Pt(8)
    
    # Protocols
    protocol_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(8.5), Inches(2.5))
    protocol_frame = protocol_box.text_frame
    protocol_frame.word_wrap = True
    
    p_title = protocol_frame.add_paragraph()
    p_title.text = "Protocols & Streaming"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = accent_red
    
    protocol_items = ["RTSP (Real-Time Streaming Protocol)", "RTP (Real-Time Transport)", "HTTP Range Requests", "1MB Chunk Streaming"]
    for item in protocol_items:
        p = protocol_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(20)
        p.font.color.rgb = white
        p.space_before = Pt(8)
    
    # Slide 5: Custom Video Player
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = dark_bg
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title = title_frame.add_paragraph()
    title.text = "🎮 Custom Video Player Controls"
    title.font.size = Pt(44)
    title.font.bold = True
    title.font.color.rgb = accent_red
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    controls = [
        "▶️ Play/Pause - Click or Space/K keys",
        "⏪ Backward 2 sec - Button or ← key",
        "⏩ Forward 2 sec - Button or → key",
        "⏮️ Skip 10 sec - Shift + arrows or J/L keys",
        "🔊 Volume Control - Slider or ↑/↓ keys",
        "🔇 Mute Toggle - Button or M key",
        "⛶ Fullscreen - Button or F key",
        "📊 Progress Bar - Clickable seeking",
        "⏱️ Time Display - Current / Total duration"
    ]
    
    for control in controls:
        p = content_frame.add_paragraph()
        p.text = control
        p.font.size = Pt(20)
        p.font.color.rgb = white
        p.space_before = Pt(10)
    
    # Slide 6: Project Structure
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = dark_bg
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title = title_frame.add_paragraph()
    title.text = "📁 Project Structure"
    title.font.size = Pt(44)
    title.font.bold = True
    title.font.color.rgb = accent_red
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    structure = [
        "web_server.py - Flask web application",
        "Server.py / Client.py - RTSP server and client",
        "RtpPacket.py - RTP packet handling",
        "assets/videos/ - Video storage (10+ videos)",
        "static/ - CSS, JavaScript files",
        "  ├── css/style.css - YouTube-style theme",
        "  └── js/player.js - Custom controls",
        "templates/ - HTML pages (index, player)",
        "requirements.txt - Python dependencies",
        "Procfile - Deployment configuration"
    ]
    
    for item in structure:
        p = content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = white if not item.startswith("  ") else gray
        p.space_before = Pt(6)
        if item.startswith("  "):
            p.level = 1
    
    # Slide 7: Deployment Journey
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = dark_bg
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title = title_frame.add_paragraph()
    title.text = "🚀 Deployment Journey"
    title.font.size = Pt(44)
    title.font.bold = True
    title.font.color.rgb = accent_red
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    steps = [
        "1️⃣ Local Development - Built and tested locally",
        "2️⃣ GitHub Repository - Pushed to version control",
        "3️⃣ Video Assets - Added 10 demo videos (70MB)",
        "4️⃣ Render.com Deployment - Connected GitHub repo",
        "5️⃣ Auto-Build Pipeline - Automated deployment",
        "6️⃣ Live Production - https://video-streaming-rtsp.onrender.com",
        "7️⃣ Documentation - Comprehensive guides created",
        "8️⃣ Live Demo Link - Added to README"
    ]
    
    for step in steps:
        p = content_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(22)
        p.font.color.rgb = white
        p.space_before = Pt(12)
    
    # Slide 8: Live Demo & Features
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = dark_bg
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title = title_frame.add_paragraph()
    title.text = "🌐 Live Demo & Access"
    title.font.size = Pt(44)
    title.font.bold = True
    title.font.color.rgb = accent_red
    
    # URL Box
    url_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1))
    url_frame = url_box.text_frame
    url = url_frame.add_paragraph()
    url.text = "https://video-streaming-rtsp.onrender.com"
    url.font.size = Pt(28)
    url.font.bold = True
    url.font.color.rgb = accent_red
    url.alignment = PP_ALIGN.CENTER
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(8.5), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    demo_features = [
        "✅ 10 High-Quality Videos Available",
        "✅ Full HD & 4K Support",
        "✅ Cross-Device Compatibility",
        "✅ Network Streaming Enabled",
        "✅ Professional UI/UX",
        "✅ Free Tier Hosting (may take 30s to wake)",
        "✅ GitHub Repository: safiullah-foragy/Video_Streaming_RTSP"
    ]
    
    for feature in demo_features:
        p = content_frame.add_paragraph()
        p.text = feature
        p.font.size = Pt(22)
        p.font.color.rgb = white
        p.space_before = Pt(12)
    
    # Slide 9: Challenges & Solutions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = dark_bg
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title = title_frame.add_paragraph()
    title.text = "💡 Challenges & Solutions"
    title.font.size = Pt(44)
    title.font.bold = True
    title.font.color.rgb = accent_red
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    challenges = [
        "🔴 Challenge: Browser autoplay restrictions",
        "   ✅ Solution: Muted autoplay + smart unmute",
        "",
        "🔴 Challenge: Custom video controls needed",
        "   ✅ Solution: Built from scratch with HTML5 API",
        "",
        "🔴 Challenge: 2-second seeking requirement",
        "   ✅ Solution: Custom buttons + keyboard shortcuts",
        "",
        "🔴 Challenge: Large video files for deployment",
        "   ✅ Solution: Selective upload (under 15MB files)",
        "",
        "🔴 Challenge: GitHub file size limits",
        "   ✅ Solution: Force-add specific videos with git -f"
    ]
    
    for item in challenges:
        p = content_frame.add_paragraph()
        p.text = item
        if item.startswith("🔴"):
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = accent_red
        elif item.startswith("   ✅"):
            p.font.size = Pt(18)
            p.font.color.rgb = white
            p.level = 1
        else:
            p.font.size = Pt(6)
        p.space_before = Pt(8)
    
    # Slide 10: Conclusion & Future
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = dark_bg
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title = title_frame.add_paragraph()
    title.text = "🎯 Summary & Future Work"
    title.font.size = Pt(44)
    title.font.bold = True
    title.font.color.rgb = accent_red
    
    # Achievements
    achieve_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(2.2))
    achieve_frame = achieve_box.text_frame
    achieve_frame.word_wrap = True
    
    a_title = achieve_frame.add_paragraph()
    a_title.text = "✨ Achievements"
    a_title.font.size = Pt(26)
    a_title.font.bold = True
    a_title.font.color.rgb = accent_red
    
    achievements = [
        "Complete video streaming solution deployed live",
        "Professional UI with custom controls",
        "Dual interface (Web + RTSP)",
        "Production-ready on Render.com"
    ]
    
    for item in achievements:
        p = achieve_frame.add_paragraph()
        p.text = "✓ " + item
        p.font.size = Pt(18)
        p.font.color.rgb = white
        p.space_before = Pt(8)
    
    # Future Enhancements
    future_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.5), Inches(2.8))
    future_frame = future_box.text_frame
    future_frame.word_wrap = True
    
    f_title = future_frame.add_paragraph()
    f_title.text = "🚀 Future Enhancements"
    f_title.font.size = Pt(26)
    f_title.font.bold = True
    f_title.font.color.rgb = accent_red
    
    future = [
        "Cloud storage integration (AWS S3, Cloudflare R2)",
        "User authentication and playlists",
        "Video upload functionality",
        "Advanced analytics and viewing stats",
        "Mobile app development"
    ]
    
    for item in future:
        p = future_frame.add_paragraph()
        p.text = "→ " + item
        p.font.size = Pt(18)
        p.font.color.rgb = gray
        p.space_before = Pt(8)
    
    # Save presentation
    prs.save('Video_Streaming_Server_Presentation.pptx')
    print("✅ Presentation created successfully!")
    print("📁 Location: Video_Streaming_Server_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
